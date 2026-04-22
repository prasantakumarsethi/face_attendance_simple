from pathlib import Path
from datetime import datetime

# ---------- PDF (ReportLab)
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.lib.units import cm
from reportlab.lib import colors

# ---------- PPT (python-pptx)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

OUT_DIR = Path(__file__).resolve().parent
PROJECT_DIR = OUT_DIR.parent

PDF_PATH = OUT_DIR / "Smart_Face_Attendance_Report.pdf"
PPT_PATH = OUT_DIR / "Smart_Face_Attendance_Presentation.pptx"

# Edit these details
META = {
    "title": "Smart Face Attendance System (Auto IN/OUT)",
    "student_name": "YOUR NAME",
    "roll_no": "YOUR ROLL NO",
    "department": "YOUR DEPARTMENT",
    "college": "YOUR COLLEGE / UNIVERSITY",
    "guide": "GUIDE NAME",
    "year": f"{datetime.now().year}",
    "tech": "Python • OpenCV (LBPH) • PySide6 • CSV Reports",
}

def draw_header(c, w, h, text):
    c.setFillColor(colors.HexColor("#0B3B74"))
    c.rect(0, h-2.2*cm, w, 2.2*cm, stroke=0, fill=1)
    c.setFillColor(colors.white)
    c.setFont("Helvetica-Bold", 16)
    c.drawString(1.5*cm, h-1.45*cm, text)

def draw_footer(c, w, h, page_no):
    c.setFillColor(colors.HexColor("#111827"))
    c.rect(0, 0, w, 1.1*cm, stroke=0, fill=1)
    c.setFillColor(colors.white)
    c.setFont("Helvetica", 9)
    c.drawString(1.5*cm, 0.4*cm, f"{META['college']}  |  {META['title']}")
    c.drawRightString(w-1.5*cm, 0.4*cm, f"Page {page_no}")

def section_title(c, x, y, text):
    c.setFillColor(colors.HexColor("#0F172A"))
    c.setFont("Helvetica-Bold", 14)
    c.drawString(x, y, text)

def paragraph(c, x, y, text, max_width=17*cm, leading=13):
    """Very simple word wrap."""
    c.setFillColor(colors.HexColor("#111827"))
    c.setFont("Helvetica", 11)
    words = text.split()
    line = ""
    lines = []
    for w in words:
        test = (line + " " + w).strip()
        if c.stringWidth(test, "Helvetica", 11) <= max_width:
            line = test
        else:
            lines.append(line)
            line = w
    if line:
        lines.append(line)
    for ln in lines:
        c.drawString(x, y, ln)
        y -= leading
    return y

def generate_pdf():
    w, h = A4
    c = canvas.Canvas(str(PDF_PATH), pagesize=A4)

    # --- Cover page
    c.setFillColor(colors.HexColor("#0B1220"))
    c.rect(0, 0, w, h, stroke=0, fill=1)

    # Accent shapes
    c.setFillColor(colors.HexColor("#2563EB"))
    c.circle(w-2.5*cm, h-2.5*cm, 2.2*cm, stroke=0, fill=1)
    c.setFillColor(colors.HexColor("#10B981"))
    c.circle(2.5*cm, 2.5*cm, 1.8*cm, stroke=0, fill=1)

    c.setFillColor(colors.white)
    c.setFont("Helvetica-Bold", 22)
    c.drawString(2*cm, h-4*cm, META["title"])

    c.setFont("Helvetica", 12)
    c.setFillColor(colors.HexColor("#CFE0FF"))
    c.drawString(2*cm, h-4.9*cm, META["tech"])

    c.setFillColor(colors.white)
    c.setFont("Helvetica-Bold", 12)
    c.drawString(2*cm, h-6.3*cm, "Submitted By:")
    c.setFont("Helvetica", 12)
    c.drawString(2*cm, h-7.0*cm, f"{META['student_name']}  ({META['roll_no']})")

    c.setFont("Helvetica-Bold", 12)
    c.drawString(2*cm, h-8.2*cm, "Department:")
    c.setFont("Helvetica", 12)
    c.drawString(2*cm, h-8.9*cm, META["department"])

    c.setFont("Helvetica-Bold", 12)
    c.drawString(2*cm, h-10.1*cm, "Guide:")
    c.setFont("Helvetica", 12)
    c.drawString(2*cm, h-10.8*cm, META["guide"])

    c.setFont("Helvetica-Bold", 12)
    c.drawString(2*cm, h-12.0*cm, "Institute:")
    c.setFont("Helvetica", 12)
    c.drawString(2*cm, h-12.7*cm, META["college"])

    c.setFont("Helvetica-Bold", 12)
    c.drawString(2*cm, h-13.9*cm, "Academic Year:")
    c.setFont("Helvetica", 12)
    c.drawString(2*cm, h-14.6*cm, META["year"])

    c.showPage()

    # --- Content pages
    page = 1

    def new_page(title):
        nonlocal page
        draw_header(c, w, h, title)
        draw_footer(c, w, h, page)
        page += 1

    # Page: Abstract
    new_page("Abstract")
    y = h - 3.2*cm
    section_title(c, 1.5*cm, y, "Abstract")
    y -= 0.9*cm
    y = paragraph(
        c, 1.5*cm, y,
        "This project presents a Smart Face Attendance System that marks IN/OUT attendance using real-time "
        "face detection and recognition. The system is touchless, reduces queues, supports multi-face capture, "
        "prevents duplicate marking, and generates exportable attendance reports. A modern UI is provided for "
        "registration, live attendance, records, editing, and exporting."
    )
    c.showPage()

    # Page: Objectives & Features
    new_page("Objectives & Key Features")
    y = h - 3.2*cm
    section_title(c, 1.5*cm, y, "Objectives")
    y -= 0.9*cm
    y = paragraph(c, 1.5*cm, y,
                  "1) Touchless attendance using camera. 2) Auto IN/OUT with configurable time rules. "
                  "3) Multi-face recognition for rapid marking. 4) Prevent duplicates & proxy attempts (blink option). "
                  "5) Export reports (daily and combined).")
    y -= 0.5*cm
    section_title(c, 1.5*cm, y, "Key Features Implemented")
    y -= 0.9*cm
    y = paragraph(c, 1.5*cm, y,
                  "• Real-time face detection (Haar Cascade)\n"
                  "• Face recognition using LBPH (OpenCV Contrib)\n"
                  "• Registration with sample capture + blur validation\n"
                  "• Auto IN/OUT logic and configurable time windows\n"
                  "• Records table with Edit/Delete\n"
                  "• Snapshot storage for audit\n"
                  "• Export to CSV")
    c.showPage()

    # Page: System Design
    new_page("System Design")
    y = h - 3.2*cm
    section_title(c, 1.5*cm, y, "Modules")
    y -= 0.9*cm
    y = paragraph(c, 1.5*cm, y,
                  "1) Registration Module: capture face samples and create dataset.\n"
                  "2) Training Module: build LBPH model and save trainer.yml.\n"
                  "3) Attendance Module: detect multiple faces, recognize, apply IN/OUT rules, save logs.\n"
                  "4) Reporting Module: export CSV, view/edit/delete logs.\n"
                  "5) Security Logging: unknown snapshots for review.")
    y -= 0.4*cm
    section_title(c, 1.5*cm, y, "Data Storage")
    y -= 0.9*cm
    y = paragraph(c, 1.5*cm, y,
                  "• dataset/<id>_<name>/ : face samples (grayscale)\n"
                  "• trainer/trainer.yml : trained LBPH model\n"
                  "• attendance_logs/attendance_YYYY-MM-DD.csv : daily logs\n"
                  "• attendance_logs/snapshots/marked : marked snapshots\n"
                  "• attendance_logs/snapshots/unknown : unknown snapshots")
    c.showPage()

    # Page: Conclusion & Future Work
    new_page("Conclusion & Future Work")
    y = h - 3.2*cm
    section_title(c, 1.5*cm, y, "Conclusion")
    y -= 0.9*cm
    y = paragraph(c, 1.5*cm, y,
                  "The system demonstrates a practical attendance workflow with real-time recognition, auto IN/OUT "
                  "handling, and modern UI. It reduces manual effort, avoids morning queues, and provides exportable "
                  "reports for institutions.")
    y -= 0.6*cm
    section_title(c, 1.5*cm, y, "Future Enhancements")
    y -= 0.9*cm
    y = paragraph(c, 1.5*cm, y,
                  "• Database (SQLite/MySQL) for large deployments\n"
                  "• Strong anti-spoofing (CNN-based liveness)\n"
                  "• ERP/School management system integration\n"
                  "• Automated absentee alerts and analytics dashboard")
    c.save()

def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(11, 18, 32)
    bg.line.fill.background()

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(37, 99, 235)
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(1))
    tf = tx.text_frame
    tf.text = META["title"]
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    sub = slide.shapes.add_textbox(Inches(0.7), Inches(2.3), Inches(12), Inches(0.8))
    sf = sub.text_frame
    sf.text = META["tech"]
    sf.paragraphs[0].font.size = Pt(18)
    sf.paragraphs[0].font.color.rgb = RGBColor(207, 224, 255)

    info = slide.shapes.add_textbox(Inches(0.7), Inches(5.0), Inches(12), Inches(1.5))
    it = info.text_frame
    it.text = f"Submitted by: {META['student_name']} ({META['roll_no']})"
    it.paragraphs[0].font.size = Pt(16)
    it.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    p = it.add_paragraph()
    p.text = f"{META['department']} | {META['college']} | Guide: {META['guide']}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(207, 224, 255)

def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tbox = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.7))
    tf = tbox.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(30)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(15, 23, 42)

    # Accent line
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.7), Inches(1.2), Inches(3.0), Inches(0.12))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(37, 99, 235)
    line.line.fill.background()

    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(12.3), Inches(5.4))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(248, 250, 252)
    box.line.color.rgb = RGBColor(226, 232, 240)

    content = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.7), Inches(5.0))
    ctf = content.text_frame
    ctf.clear()

    for i, b in enumerate(bullets):
        p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(30, 41, 59)

def generate_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    add_bullets_slide(prs, "Problem Statement", [
        "Manual attendance is time-consuming and prone to proxy marking.",
        "Card/Fingerprint systems need contact and cause queues.",
        "Institutions need fast, touchless, reliable attendance with reports."
    ])

    add_bullets_slide(prs, "Objectives", [
        "Real-time face-based attendance with Auto IN/OUT.",
        "Touchless workflow to reduce queues and improve hygiene.",
        "Configurable time rules for IN/OUT and Late marking.",
        "Exportable reports (daily and combined)."
    ])

    add_bullets_slide(prs, "System Architecture", [
        "Registration: capture face samples → dataset",
        "Training: LBPH model → trainer.yml",
        "Attendance: detect faces → recognize → apply rules → store CSV",
        "Records: view/edit/delete → export reports"
    ])

    add_bullets_slide(prs, "Key Features", [
        "Multi-face marking (no morning queue)",
        "Auto IN/OUT logic and duplicate prevention",
        "Optional blink-based anti-proxy (demo liveness)",
        "Unknown snapshot logging for security review",
        "Modern UI with dashboard and exports"
    ])

    add_bullets_slide(prs, "Results / Demo Plan", [
        "Register student → Train model",
        "Start attendance → walk in front of camera → popup confirms mark",
        "View records → export CSV",
        "Show IN then OUT after allowed time"
    ])

    add_bullets_slide(prs, "Future Enhancements", [
        "Database (SQLite/MySQL) for larger scale",
        "Advanced anti-spoofing (CNN liveness)",
        "ERP/School Management integration",
        "Absentee alerts and analytics dashboard"
    ])

    prs.save(str(PPT_PATH))

def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    generate_pdf()
    generate_ppt()
    print("Generated:")
    print(" -", PDF_PATH)
    print(" -", PPT_PATH)

if __name__ == "__main__":
    main()